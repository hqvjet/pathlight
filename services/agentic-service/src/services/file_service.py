import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from io import BytesIO
from fastapi import UploadFile


def extract_content_with_tags(file: UploadFile, extension: str) -> str:
    """
    Extract content from uploaded files and add structured tags.
    
    Args:
        file: The uploaded file object
        extension: File extension (pdf, docx, pptx)
        
    Returns:
        str: Tagged content extracted from the file
    """
    tagged_content = ""

    if extension == "pdf":
        # Reset file pointer to beginning
        file.file.seek(0)
        pdf_document = fitz.open(stream=file.file.read(), filetype="pdf")
        for page_number, page in enumerate(pdf_document, start=1):
            page_text = page.get_text()
            tagged_content += f'<page number="{page_number}">{page_text}</page>\n'
        pdf_document.close()

    elif extension == "docx":
        # Reset file pointer to beginning
        file.file.seek(0)
        doc = Document(file.file)
        
        # Process paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.style.name.startswith("Heading"):
                level = paragraph.style.name.replace("Heading", "h")
                tagged_content += f'<heading level="{level}">{paragraph.text}</heading>\n'
            else:
                tagged_content += f'<paragraph>{paragraph.text}</paragraph>\n'
        
        # Process tables
        for table in doc.tables:
            table_rows = ""
            for row in table.rows:
                row_cells = "".join([f'<cell>{cell.text}</cell>' for cell in row.cells])
                table_rows += f'<row>{row_cells}</row>'
            tagged_content += f'<table>{table_rows}</table>\n'

    elif extension == "pptx":
        # Reset file pointer to beginning
        file.file.seek(0)
        presentation = Presentation(file.file)
        
        for slide_number, slide in enumerate(presentation.slides, start=1):
            slide_title = ""
            slide_content = ""
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if hasattr(slide.shapes, 'title') and shape == slide.shapes.title:
                        slide_title = f'<slide_title>{shape.text}</slide_title>'
                    else:
                        slide_content += f'<slide_content>{shape.text}</slide_content>'
            
            tagged_content += f'<slide number="{slide_number}">{slide_title}{slide_content}</slide>\n'

    return tagged_content