import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from fastapi import UploadFile


def extract_content_with_tags(file: UploadFile, extension: str):
    tagged_content = ""

    if extension == "pdf":
        pdf_document = fitz.open(stream=file.file.read(), filetype="pdf")
        for page_number, page in enumerate(pdf_document, start=1):
            page_text = page.get_text()
            tagged_content += f'<page number="{page_number}">{page_text}</page>'

    elif extension == "docx":
        doc = Document(file.file)
        for paragraph in doc.paragraphs:
            if paragraph.style.name.startswith("Heading"):
                level = paragraph.style.name.replace("Heading", "h")
                tagged_content += f'<heading level="{level}">{paragraph.text}</heading>'
            else:
                tagged_content += f'<paragraph>{paragraph.text}</paragraph>'
        for table in doc.tables:
            table_rows = ""
            for row in table.rows:
                row_cells = "".join([f'<cell>{cell.text}</cell>' for cell in row.cells])
                table_rows += f'<row>{row_cells}</row>'
            tagged_content += f'<table>{table_rows}</table>'

    elif extension == "pptx":
        presentation = Presentation(file.file)
        for slide_number, slide in enumerate(presentation.slides, start=1):
            slide_title = ""
            slide_content = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text.strip():
                        if shape == slide.shapes.title:
                            slide_title = f'<slide_title>{shape.text}</slide_title>'
                        else:
                            slide_content += f'<slide_content>{shape.text}</slide_content>'
            tagged_content += f'<slide number="{slide_number}">{slide_title}{slide_content}</slide>'

    return tagged_content